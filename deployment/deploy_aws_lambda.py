import json
import boto3
from firecrawl import FirecrawlApp
from pptx import Presentation
import uuid
import os

# Configuration
FIRECRAWL_API_KEY = os.environ.get('FIRECRAWL_API_KEY', 'your-api-key-here')
S3_BUCKET_NAME = os.getenv("S3_BUCKET_NAME")

def lambda_handler(event, context):
    try:
        # Parse request
        body = json.loads(event['body'])
        description = body['description']
        urls = body['urls']
        
        # Crawl URLs
        firecrawl = FirecrawlApp(api_key=FIRECRAWL_API_KEY)
        content_data = []
        
        for url in urls:
            try:
                result = firecrawl.scrape_url(url)
                content_data.append({
                    "url": url,
                    "content": result.get("content", "")[:2000]
                })
            except:
                continue
        
        # Generate slides with Bedrock
        bedrock = boto3.client('bedrock-runtime')
        
        prompt = f"""Create a PowerPoint presentation based on:
Description: {description}

Content from URLs:
{json.dumps(content_data, indent=2)}

Generate 5-6 slides with clear headings and bullet points.
Format as JSON: {{"slides": [{{"title": "Title", "content": ["Point 1", "Point 2"]}}]}}"""

        response = bedrock.invoke_model(
            modelId='us.anthropic.claude-3-5-sonnet-20241022-v2:0',
            body=json.dumps({
                "anthropic_version": "bedrock-2023-05-31",
                "max_tokens": 2000,
                "messages": [{"role": "user", "content": prompt}]
            })
        )
        
        result = json.loads(response['body'].read())
        slides_json = json.loads(result['content'][0]['text'])
        
        # Create PowerPoint
        prs = Presentation()
        
        for slide_data in slides_json['slides']:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = slide_data['title']
            content.text = '\n'.join([f"• {point}" for point in slide_data['content']])
        
        # Save to S3
        file_name = f"presentation_{uuid.uuid4().hex[:8]}.pptx"
        local_path = f"/tmp/{file_name}"
        
        prs.save(local_path)
        
        s3 = boto3.client('s3')
        s3.upload_file(local_path, S3_BUCKET_NAME, file_name)
        
        download_url = s3.generate_presigned_url(
            'get_object',
            Params={'Bucket': S3_BUCKET_NAME, 'Key': file_name},
            ExpiresIn=3600
        )
        
        return {
            'statusCode': 200,
            'headers': {'Content-Type': 'application/json'},
            'body': json.dumps({'download_url': download_url})
        }
        
    except Exception as e:
        return {
            'statusCode': 500,
            'body': json.dumps({'error': str(e)})
        }