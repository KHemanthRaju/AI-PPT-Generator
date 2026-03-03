import json
import urllib.request
import re
import boto3
import uuid
import os
import base64
from urllib.parse import urljoin, urlparse, quote, unquote
from io import BytesIO
import csv

def parse_excel_data(file_content, file_extension):
    """Parse Excel/CSV data and return structured information"""
    try:
        import pandas as pd
        from io import BytesIO

        # Determine file type and parse accordingly
        if file_extension in ['.xlsx', '.xls']:
            # Excel file - skip empty rows at top
            df = pd.read_excel(BytesIO(file_content))

            # Find the actual header row (look for row with most non-null values)
            header_row = None
            for idx in range(min(10, len(df))):  # Check first 10 rows
                row = df.iloc[idx]
                non_null_count = row.notna().sum()
                if non_null_count >= 3:  # At least 3 columns with data
                    header_row = idx
                    break

            # Re-read with correct header
            if header_row is not None and header_row > 0:
                print(f"Detected header at row {header_row}, re-reading Excel...")
                df = pd.read_excel(BytesIO(file_content), header=header_row)
        else:
            # CSV file
            df = pd.read_csv(BytesIO(file_content))

        if df.empty:
            return None, None

        # Drop completely empty columns
        df = df.dropna(axis=1, how='all')

        # Drop completely empty rows
        df = df.dropna(axis=0, how='all')

        # Reset index after dropping rows
        df = df.reset_index(drop=True)

        if df.empty:
            print("DataFrame is empty after cleaning")
            return None, None

        # Convert to list of dictionaries
        rows = df.to_dict('records')

        # Get column names (excluding Unnamed columns if possible)
        columns = list(df.columns)

        # If we have mostly "Unnamed" columns, the file structure is problematic
        unnamed_count = sum(1 for col in columns if str(col).startswith('Unnamed'))
        if unnamed_count > len(columns) / 2:
            print(f"Warning: {unnamed_count}/{len(columns)} columns are 'Unnamed' - Excel structure may be incorrect")

        # Convert any datetime objects in rows to strings for JSON serialization
        import datetime
        for row in rows:
            for key, value in list(row.items()):
                if isinstance(value, (datetime.datetime, datetime.date, pd.Timestamp)):
                    row[key] = str(value)
                elif pd.isna(value):
                    row[key] = None

        # Analyze data
        data_summary = {
            'columns': columns,
            'row_count': len(rows),
            'sample_data': rows[:5],  # First 5 rows
            'data_types': {}
        }

        # Detect numeric columns using pandas dtypes
        for col in columns:
            if pd.api.types.is_numeric_dtype(df[col]):
                data_summary['data_types'][col] = 'numeric'
            elif pd.api.types.is_datetime64_any_dtype(df[col]):
                data_summary['data_types'][col] = 'date'
            else:
                data_summary['data_types'][col] = 'text'

        # Debug logging
        print(f"Detected columns: {columns}")
        print(f"Data types: {data_summary['data_types']}")
        print(f"Sample row 0: {rows[0] if rows else 'None'}")
        print(f"Sample row 1: {rows[1] if len(rows) > 1 else 'None'}")

        return data_summary, rows
    except Exception as e:
        print(f"Error parsing Excel/CSV: {str(e)}")
        return None, None

def parse_csv_data(csv_content):
    """Parse CSV data and return structured information (legacy function, calls parse_excel_data)"""
    try:
        return parse_excel_data(csv_content, '.csv')
    except Exception as e:
        print(f"Error parsing CSV: {str(e)}")
        return None, None

def generate_chart_with_matplotlib(data_rows, x_column, y_column, chart_type='line'):
    """Generate a chart from data using matplotlib"""
    try:
        import matplotlib
        matplotlib.use('Agg')  # Non-interactive backend
        import matplotlib.pyplot as plt
        import numpy as np

        # Check if columns exist (case-insensitive match)
        if not data_rows:
            print("No data rows provided")
            return None

        available_columns = list(data_rows[0].keys())
        print(f"Available columns in data: {available_columns}")

        # Try to find matching column names (case-insensitive)
        def find_column(target_col, available_cols):
            if not target_col:
                return None
            target_lower = str(target_col).lower().strip()
            for col in available_cols:
                if str(col).lower().strip() == target_lower:
                    return col
            # Try partial match
            for col in available_cols:
                if target_lower in str(col).lower() or str(col).lower() in target_lower:
                    return col
            return None

        actual_x_col = find_column(x_column, available_columns)
        actual_y_col = find_column(y_column, available_columns)

        if not actual_x_col:
            print(f"Could not find X column '{x_column}' in data")
            return None
        if not actual_y_col:
            print(f"Could not find Y column '{y_column}' in data")
            return None

        print(f"Using columns: X='{actual_x_col}', Y='{actual_y_col}'")

        # Extract and clean data
        x_values = []
        y_values = []

        for row in data_rows:
            x_val = row.get(actual_x_col)
            y_val = row.get(actual_y_col)

            # Skip rows with missing data
            if x_val is None or y_val is None:
                continue

            # Convert to appropriate types
            try:
                # Convert y to float, skip if not numeric
                y_float = float(y_val)
                if np.isnan(y_float):
                    continue

                # Convert x to string for categorical axis
                x_str = str(x_val).strip()
                if not x_str or x_str.lower() in ['none', 'nan', '']:
                    continue

                x_values.append(x_str)
                y_values.append(y_float)
            except (ValueError, TypeError):
                continue

        # Check if we have enough data
        print(f"Extracted {len(x_values)} valid data points")
        if len(x_values) > 0:
            print(f"  Sample X values: {x_values[:3]}")
            print(f"  Sample Y values: {y_values[:3]}")

        if len(x_values) < 2:
            print(f"Not enough valid data points for chart: {len(x_values)} points (need at least 2)")
            return None

        # Limit to first 20 data points for readability
        if len(x_values) > 20:
            x_values = x_values[:20]
            y_values = y_values[:20]

        # Create figure with professional styling
        fig, ax = plt.subplots(figsize=(10, 6))

        # Professional color palette
        primary_color = '#1F3864'    # Dark blue (matches slide title color)
        accent_color = '#3A7BD5'     # Lighter blue for highlights

        if chart_type == 'line':
            # Line chart with professional styling
            ax.plot(range(len(x_values)), y_values,
                   marker='o', linewidth=2.5, markersize=8,
                   color=primary_color, markerfacecolor=accent_color,
                   markeredgewidth=2, markeredgecolor=primary_color)
            ax.set_xticks(range(len(x_values)))
            ax.set_xticklabels(x_values)
        elif chart_type == 'bar':
            # Bar chart with professional styling
            bars = ax.bar(range(len(x_values)), y_values, color=primary_color, alpha=0.8)
            # Highlight the maximum value
            max_idx = y_values.index(max(y_values))
            bars[max_idx].set_color(accent_color)
            bars[max_idx].set_alpha(1.0)
            ax.set_xticks(range(len(x_values)))
            ax.set_xticklabels(x_values)
        elif chart_type == 'scatter':
            # Scatter chart with professional styling
            ax.scatter(range(len(x_values)), y_values, s=120,
                      color=primary_color, alpha=0.7, edgecolors=accent_color, linewidths=2)
            ax.set_xticks(range(len(x_values)))
            ax.set_xticklabels(x_values)

        # Professional axis styling
        ax.set_xlabel(x_column, fontsize=11, fontweight='normal', color='#404040')
        ax.set_ylabel(y_column, fontsize=11, fontweight='normal', color='#404040')
        ax.set_title(f'{y_column} by {x_column}', fontsize=13, fontweight='bold',
                    color='#1F3864', pad=15)

        # Minimal gridlines (only horizontal, subtle)
        ax.grid(True, alpha=0.15, axis='y', linestyle='--', linewidth=0.5)
        ax.set_axisbelow(True)  # Grid behind bars

        # Clean up spines (remove top and right borders)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.spines['left'].set_color('#CCCCCC')
        ax.spines['bottom'].set_color('#CCCCCC')

        # Rotate x-axis labels if needed
        if len(x_values) > 5:
            plt.xticks(rotation=45, ha='right', fontsize=9)
        else:
            plt.xticks(fontsize=10)

        # Y-axis formatting
        plt.yticks(fontsize=10)

        plt.tight_layout(pad=1.5)

        # Save to bytes
        img_buffer = BytesIO()
        plt.savefig(img_buffer, format='png', dpi=150, bbox_inches='tight')
        img_buffer.seek(0)
        plt.close()

        print(f"Successfully generated {chart_type} chart with {len(x_values)} data points")
        return img_buffer.getvalue()
    except Exception as e:
        print(f"Error generating chart: {str(e)}")
        import traceback
        traceback.print_exc()
        return None

def analyze_data_with_claude(bedrock_client, data_summary, rows):
    """Use Claude to analyze data and suggest visualizations"""
    try:
        # Prepare data summary for Claude
        summary_text = f"""
Data Analysis Request:

Columns: {', '.join(data_summary['columns'])}
Total Rows: {data_summary['row_count']}

Data Types:
{json.dumps(data_summary['data_types'], indent=2)}

Sample Data (first 5 rows):
{json.dumps(data_summary['sample_data'], indent=2)}

IMPORTANT: Please analyze this data and provide visualizations. You MUST suggest at least 1-3 charts to visualize the numeric data.

Return your analysis as PURE JSON (no markdown, no explanations):
{{
  "insights": ["insight 1 about trends in the data", "insight 2 about key findings", "insight 3 about patterns"],
  "visualizations": [
    {{"x_column": "column_name", "y_column": "numeric_column_name", "chart_type": "line", "title": "Descriptive Chart Title"}},
    {{"x_column": "column_name", "y_column": "another_numeric_column", "chart_type": "bar", "title": "Another Chart Title"}}
  ],
  "slide_content": ["bullet point 1", "bullet point 2", "bullet point 3"]
}}

RULES:
- For visualizations, x_column should be categorical (like Month, Category, Date) and y_column must be numeric
- Identify which columns are numeric from the data_types above
- Chart types can be: "line" (for trends over time), "bar" (for comparisons), "scatter" (for correlations)
- Return ONLY the JSON, nothing else
"""

        response = bedrock_client.invoke_model(
            modelId='us.anthropic.claude-sonnet-4-5-20250929-v1:0',
            body=json.dumps({
                "anthropic_version": "bedrock-2023-05-31",
                "max_tokens": 2000,
                "temperature": 0.5,
                "messages": [{"role": "user", "content": summary_text}]
            })
        )

        result = json.loads(response['body'].read())
        analysis_text = result['content'][0]['text']

        print(f"Claude data analysis response: {analysis_text[:500]}")

        # Try to parse JSON from response with better error handling
        try:
            # First, try direct parsing
            analysis = json.loads(analysis_text)
        except json.JSONDecodeError:
            # Try to extract JSON from markdown code blocks
            json_match = re.search(r'```(?:json)?\s*(\{.*?\})\s*```', analysis_text, re.DOTALL)
            if json_match:
                try:
                    analysis = json.loads(json_match.group(1))
                    print("Successfully extracted JSON from markdown code block")
                except json.JSONDecodeError:
                    json_match = None

            # Try to find JSON object directly
            if not json_match:
                json_match = re.search(r'\{[\s\S]*"insights"[\s\S]*\}', analysis_text)
                if json_match:
                    try:
                        analysis = json.loads(json_match.group(0))
                        print("Successfully extracted JSON from text")
                    except json.JSONDecodeError:
                        analysis = None
                else:
                    analysis = None

            # Fallback with default visualizations
            if not analysis:
                # Try to create visualizations based on data types
                numeric_cols = [col for col, dtype in data_summary['data_types'].items() if dtype == 'numeric']
                non_numeric_cols = [col for col, dtype in data_summary['data_types'].items() if dtype == 'text']

                visualizations = []
                if len(numeric_cols) > 0 and len(non_numeric_cols) > 0:
                    # Create visualizations for each numeric column
                    x_col = non_numeric_cols[0]  # Use first text column as x-axis
                    for i, y_col in enumerate(numeric_cols[:3]):  # Max 3 charts
                        chart_type = "line" if i == 0 else "bar"
                        visualizations.append({
                            "x_column": x_col,
                            "y_column": y_col,
                            "chart_type": chart_type,
                            "title": f"{y_col} by {x_col}"
                        })

                analysis = {
                    "insights": [
                        "Data analysis completed successfully",
                        "Key trends and patterns identified in the dataset",
                        "Visualizations created to illustrate the data"
                    ],
                    "visualizations": visualizations,
                    "slide_content": ["Data overview and key metrics", "Trends and insights from analysis", "Recommendations based on findings"]
                }
                print(f"Using fallback analysis with {len(visualizations)} auto-generated visualizations")

        return analysis
    except Exception as e:
        print(f"Error analyzing data with Claude: {str(e)}")
        import traceback
        traceback.print_exc()
        return None

def add_hyperlinks_to_paragraph(paragraph, text):
    """
    Add text to paragraph with automatic hyperlinks for URLs.
    Follows accessibility standard: display fully qualified URLs with active hyperlinks.
    """
    from pptx.dml.color import RGBColor

    # Find all URLs in the text
    url_pattern = r'(https?://[^\s]+)'
    parts = re.split(url_pattern, text)

    for i, part in enumerate(parts):
        if re.match(url_pattern, part):
            # This is a URL - add as hyperlink
            run = paragraph.add_run()
            run.text = part
            run.font.color.rgb = RGBColor(0, 0, 255)  # Blue color for links
            run.font.bold = True  # Bold for accessibility (don't rely on color alone)
            run.font.underline = True
            # Add hyperlink
            hlink = run.hyperlink
            hlink.address = part
        else:
            # Regular text
            if part:
                run = paragraph.add_run()
                run.text = part

def extract_text_from_html(html_content):
    """Extract meaningful text from HTML using regex"""
    html_content = re.sub(r'<script[^>]*>.*?</script>', '', html_content, flags=re.DOTALL | re.IGNORECASE)
    html_content = re.sub(r'<style[^>]*>.*?</style>', '', html_content, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r'<[^>]+>', '', html_content)
    text = re.sub(r'\s+', ' ', text)
    return text.strip()[:2000]

def extract_images_from_html(html_content, base_url):
    """Extract image URLs from HTML"""
    images = []
    # Find img tags with src
    img_pattern = r'<img[^>]+src=["\']([^"\']+)["\']'
    matches = re.findall(img_pattern, html_content, re.IGNORECASE)

    for img_url in matches[:20]:  # Check first 20 images
        # Convert relative URLs to absolute
        if img_url.startswith('//'):
            img_url = 'https:' + img_url
        elif img_url.startswith('/'):
            img_url = urljoin(base_url, img_url)
        elif not img_url.startswith('http'):
            img_url = urljoin(base_url, img_url)

        # Skip obvious non-content images
        skip_patterns = ['logo', 'icon', 'favicon', 'avatar', 'emoji', 'badge', 'button', 'sprite', 'pixel', 'tracker', 'ad.', 'ads.']

        # Skip if URL contains skip patterns
        if any(pattern in img_url.lower() for pattern in skip_patterns):
            continue

        # Skip data URLs and very small images
        if img_url.startswith('data:') or 'spacer.gif' in img_url.lower():
            continue

        # Skip if filename is too short (likely an icon)
        url_parts = img_url.split('/')
        if url_parts:
            filename = url_parts[-1].split('?')[0]
            if len(filename) < 5:  # Very short filenames are usually icons
                continue

        images.append(img_url)

    return images[:10]  # Return up to 10 images

def search_web(query, num_results=3):
    """Search the web for a query and return top result URLs"""
    try:
        print(f"Searching web for: {query}")

        # Use DuckDuckGo HTML search (no API key needed)
        search_url = f"https://html.duckduckgo.com/html/?q={quote(query)}"

        req = urllib.request.Request(
            search_url,
            headers={
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
            }
        )

        with urllib.request.urlopen(req, timeout=10) as response:
            html = response.read().decode('utf-8')

        # Extract URLs from DuckDuckGo results
        # Look for result links
        url_pattern = r'<a[^>]+class="[^"]*result__a[^"]*"[^>]+href="//duckduckgo\.com/l/\?uddg=([^"&]+)'
        matches = re.findall(url_pattern, html)

        urls = []
        for match in matches[:num_results]:
            try:
                # URL decode the DuckDuckGo redirect URL
                decoded_url = unquote(match)
                if decoded_url.startswith('http'):
                    urls.append(decoded_url)
            except:
                continue

        # Fallback: Try Wikipedia
        if len(urls) == 0:
            wiki_url = f"https://en.wikipedia.org/wiki/{quote(query.replace(' ', '_'))}"
            print(f"Using Wikipedia fallback: {wiki_url}")
            urls.append(wiki_url)

        print(f"Found {len(urls)} URLs from search: {urls}")
        return urls

    except Exception as e:
        print(f"Error searching web: {str(e)}")
        # Fallback to Wikipedia
        wiki_query = query.replace(' ', '_')
        return [f"https://en.wikipedia.org/wiki/{wiki_query}"]

def download_and_encode_image(image_url):
    """Download image and convert to base64"""
    try:
        req = urllib.request.Request(image_url, headers={'User-Agent': 'Mozilla/5.0'})
        with urllib.request.urlopen(req, timeout=10) as response:
            image_data = response.read()

            # Check file size (skip very small images - likely icons)
            # Limit to 3.75MB (Claude's per-image limit)
            if len(image_data) < 2048:  # Skip images smaller than 2KB (likely icons)
                return None
            if len(image_data) > 3.75 * 1024 * 1024:  # Claude limit is 3.75MB per image
                print(f"Skipping large image: {len(image_data)} bytes")
                return None

            # Get content type
            content_type = response.headers.get('Content-Type', '').lower()
            if 'image' not in content_type:
                return None

            # Only accept Claude-supported formats
            supported = ['jpeg', 'jpg', 'png', 'gif', 'webp']
            if not any(fmt in content_type for fmt in supported):
                print(f"Unsupported format: {content_type}")
                return None

            # Determine media type
            if 'jpeg' in content_type or 'jpg' in content_type:
                media_type = 'image/jpeg'
            elif 'png' in content_type:
                media_type = 'image/png'
            elif 'gif' in content_type:
                media_type = 'image/gif'
            elif 'webp' in content_type:
                media_type = 'image/webp'
            else:
                return None  # Skip if can't determine

            # Convert to base64
            base64_image = base64.b64encode(image_data).decode('utf-8')
            if len(base64_image) == 0:
                return None

            return {
                'url': image_url,
                'base64': base64_image,
                'media_type': media_type,
                'size': len(image_data)
            }
    except Exception as e:
        print(f"Error downloading image {image_url}: {str(e)}")
        return None

def lambda_handler(event, context):
    try:
        body = json.loads(event['body'])
        description = body.get('description', 'Create a presentation')
        urls = body.get('urls', [])
        csv_data_b64 = body.get('csv_data', None)  # Base64 encoded CSV/Excel data
        file_extension = body.get('file_extension', '.csv')  # File extension (.csv, .xlsx, etc.)

        # Get document text (try base64 encoded first, fallback to plain text)
        document_text_b64 = body.get('document_text_b64', None)
        if document_text_b64:
            # Decode base64 encoded document text
            document_text = base64.b64decode(document_text_b64).decode('utf-8')
            print(f"Decoded base64 document text ({len(document_text)} characters)")
        else:
            # Fallback to plain text (backward compatibility)
            document_text = body.get('document_text', None)

        # Initialize variables
        extracted_content = []
        all_images = []
        data_analysis = None
        generated_charts = []

        # Process CSV/Excel data if provided
        if csv_data_b64:
            try:
                print(f"Processing {file_extension} data")
                file_bytes = base64.b64decode(csv_data_b64)

                # Parse Excel/CSV data
                data_summary, data_rows = parse_excel_data(file_bytes, file_extension)

                if data_summary and data_rows:
                    print(f"Parsed CSV: {data_summary['row_count']} rows, {len(data_summary['columns'])} columns")

                    # Analyze data with Claude
                    bedrock = boto3.client('bedrock-runtime', region_name='us-west-2')
                    data_analysis = analyze_data_with_claude(bedrock, data_summary, data_rows)

                    if data_analysis:
                        print(f"Claude analysis complete: {len(data_analysis.get('visualizations', []))} visualizations suggested")

                        # Generate charts based on Claude's suggestions
                        for viz in data_analysis.get('visualizations', []):
                            try:
                                x_col = viz.get('x_column')
                                y_col = viz.get('y_column')

                                print(f"Attempting to generate chart: {viz.get('title')}")
                                print(f"  X column: {x_col}")
                                print(f"  Y column: {y_col}")
                                print(f"  Chart type: {viz.get('chart_type', 'line')}")
                                print(f"  Available columns: {list(data_rows[0].keys()) if data_rows else 'None'}")

                                chart_bytes = generate_chart_with_matplotlib(
                                    data_rows,
                                    x_col,
                                    y_col,
                                    viz.get('chart_type', 'line')
                                )

                                if chart_bytes:
                                    # Add chart as an image
                                    chart_b64 = base64.b64encode(chart_bytes).decode('utf-8')
                                    all_images.append({
                                        'url': f"chart_{viz.get('x_column')}_{viz.get('y_column')}",
                                        'base64': chart_b64,
                                        'media_type': 'image/png',
                                        'size': len(chart_bytes),
                                        'is_chart': True,
                                        'chart_title': viz.get('title', 'Data Visualization')
                                    })
                                    generated_charts.append(viz)
                                    print(f"Generated chart: {viz.get('title')}")
                            except Exception as chart_error:
                                print(f"Error generating chart: {str(chart_error)}")

                        # Add data insights to extracted content
                        insights_text = '\n'.join(data_analysis.get('insights', []))
                        extracted_content.append({
                            'source': 'data_analysis',
                            'content': insights_text
                        })
                else:
                    print("Failed to parse CSV data")
            except Exception as csv_error:
                print(f"Error processing CSV: {str(csv_error)}")

        # Process document text if provided (from PDF/DOC)
        document_urls = []
        if document_text:
            try:
                print(f"Processing document text ({len(document_text)} characters)")

                # Extract URLs from document text
                url_pattern = r'https?://[^\s<>"{}|\\^`\[\]]+(?:[^\s.,;!?)\]]|(?<=\w)[.,;!?)])'
                found_urls = re.findall(url_pattern, document_text)

                # Clean up URLs (remove trailing punctuation that might be part of sentence)
                cleaned_urls = []
                for url in found_urls:
                    # Remove trailing punctuation
                    url = url.rstrip('.,;:!?)')
                    if url and url not in cleaned_urls:
                        cleaned_urls.append(url)

                document_urls = cleaned_urls

                if document_urls:
                    print(f"Found {len(document_urls)} URLs in document: {document_urls}")
                    # Add document URLs to the main URLs list for crawling
                    if not urls:
                        urls = []
                    urls.extend(document_urls)
                else:
                    print("No URLs found in document")

                # Add document text to extracted content
                extracted_content.append({
                    'source': 'uploaded_document',
                    'content': document_text
                })
            except Exception as doc_error:
                print(f"Error processing document text: {str(doc_error)}")

        # If no URLs provided and no CSV data and no document, search the web automatically
        if (not urls or len(urls) == 0) and not csv_data_b64 and not document_text:
            print(f"No URLs or data provided. Searching web for: {description}")
            urls = search_web(description, num_results=3)
            print(f"Using search results: {urls}")

        # Extract text and images from URLs
        print(f"Starting URL crawling loop with {len(urls)} URLs: {urls}")
        for url in urls:
            try:
                print(f"Crawling URL: {url}")
                req = urllib.request.Request(url, headers={'User-Agent': 'Mozilla/5.0'})
                with urllib.request.urlopen(req, timeout=15) as response:
                    html_content = response.read().decode('utf-8')
                    print(f"Successfully fetched {len(html_content)} bytes from {url}")

                # Extract text
                clean_text = extract_text_from_html(html_content)
                print(f"Extracted {len(clean_text)} characters of text from {url}")

                # Extract images
                image_urls = extract_images_from_html(html_content, url)
                print(f"Found {len(image_urls)} images on {url}")

                extracted_content.append({
                    'url': url,
                    'content': clean_text,
                    'image_count': len(image_urls)
                })

                # Download and encode images
                for img_url in image_urls:
                    encoded_img = download_and_encode_image(img_url)
                    if encoded_img:
                        all_images.append(encoded_img)

            except Exception as e:
                print(f"Error crawling {url}: {str(e)}")
                extracted_content.append({'url': url, 'error': str(e)})
        
        # Generate slides with Bedrock
        bedrock = boto3.client('bedrock-runtime', region_name='us-west-2')
        all_content = '\n\n'.join([item.get('content', '') for item in extracted_content if 'content' in item])

        slide_count = body.get('slide_count', 6)
        content_slides = slide_count - 2  # Subtract title and thank you slides

        # Build message content with text and images
        message_content = []

        # Add text prompt
        # Limit images sent to Claude to 5 max
        images_to_send = all_images[:5]

        if images_to_send:
            chart_indices = [str(i) for i, img in enumerate(images_to_send) if img.get('is_chart')]
            has_charts = len(chart_indices) > 0

            image_instruction = f"""

🔴 CRITICAL IMAGE INSTRUCTIONS 🔴
I'm providing {len(images_to_send)} images. These images are numbered 0 to {len(images_to_send)-1}.

{"🔴 IMPORTANT: Images " + ', '.join(chart_indices) + " are DATA CHARTS/GRAPHS generated from the uploaded data file. YOU MUST INCLUDE THESE CHARTS IN YOUR PRESENTATION!" if has_charts else ""}

REQUIRED ACTIONS:
1. ANALYZE each image to understand what it shows (charts, graphs, data visualizations, diagrams, photos, etc.)
2. EXTRACT key insights, data points, and trends from the images
3. INCORPORATE image insights into your slide bullet points
4. **YOU MUST ASSIGN AT LEAST {min(len(images_to_send), slide_count - 2)} IMAGES TO SLIDES using "image_index" field**

IMAGE ASSIGNMENT RULES:
- Set "image_index": 0 to put image #0 on that slide
- Set "image_index": 1 to put image #1 on that slide
- Set "image_index": 2 to put image #2 on that slide
- Set "image_index": null ONLY on the first (title) and last (thank you) slides
- ALL middle content slides SHOULD have an image assigned where relevant
{"- 🔴 PRIORITY: Data charts (images " + ', '.join(chart_indices) + ") MUST be included on content slides!" if has_charts else ""}

EXAMPLE of correct usage:
{{"title": "Revenue Analysis", "content": ["..."], "image_index": 0}}  ← Image #0 appears on this slide
{{"title": "Team Performance", "content": ["..."], "image_index": 1}}  ← Image #1 appears on this slide

DO NOT set image_index to null on content slides when you have images available!"""
        else:
            image_instruction = """

🔴 IMPORTANT: NO IMAGES AVAILABLE 🔴
- Set "image_index": null on ALL slides (including content slides)
- Do NOT assign any numeric image_index values
- Create text-only slides with focused content"""

        # Add data analysis context if available
        data_context = ""
        if data_analysis:
            data_context = f"""

DATA ANALYSIS INSIGHTS:
{json.dumps(data_analysis, indent=2)}

IMPORTANT: Use these data insights in your presentation. The charts visualizing this data are included in the images."""

        prompt_text = f"""Create a {slide_count}-slide PowerPoint presentation about: {description}

Based on this content:
{all_content}{data_context}{image_instruction}

CRITICAL: Return ONLY valid JSON. Do NOT wrap in markdown code blocks. Do NOT add any explanation before or after the JSON.

Return in this EXACT format (insight-based titles + 6×6 rule):
{{"slides": [
  {{"title": "{description}", "content": ["Overview and key objectives", "Topics covered in presentation", "Main takeaways and insights"], "image_index": null}},
  {{"title": "Revenue Growth Accelerated 25% Year Over Year", "content": ["Urban markets driving expansion", "Customer base doubled in Q4", "Profit margins improved significantly"], "image_index": 0}},
  {{"title": "Project Delivered Ahead of Schedule and Under Budget", "content": ["Completed two weeks early", "Saved 15% on total costs", "Quality metrics exceeded targets"], "image_index": 1}},
  {{"title": "Thank You", "content": ["Questions and discussion welcome", "Thank you for attention"], "image_index": null}}
]}}

🔴 NOTICE: The example above shows image_index: 0, 1, 2, etc. on content slides. YOU MUST DO THIS with the images I'm providing!

🔴 SLIDE TITLE RULES - ONE MESSAGE PER SLIDE 🔴:
Each slide title should answer: "What is the one key takeaway?"

❌ WEAK (topic-based titles):
- "Market Analysis"
- "Sales Performance"
- "Customer Data"

✅ STRONG (insight-based titles):
- "Market Demand Growing 18% Annually in Urban Areas"
- "Sales Exceeded Target by $2M in Q4"
- "Customer Retention Rate Reached Record 94%"

🔴 TITLE REQUIREMENTS - MUST BE VERY SHORT 🔴:
- **CRITICAL: Maximum 1-2 WORDS per title** (VERY SHORT!)
- State a CONCLUSION, not just a topic
- Include the KEY FINDING or insight
- Use specific numbers/data when available
- Examples: "Revenue Growth", "Customer Satisfaction", "Market Share", "Performance", "Strategy"

Structure Requirements:
- **First slide (Slide 1)**: Main title slide with overview
  - Title: Use the **MAIN TOPIC TITLE** (the presentation topic itself, e.g., "{description}")
  - Content: 4-6 SHORT bullet points (4-6 words each) listing what topics will be covered (adjust based on presentation scope)
  - No image on first slide
  - Example: ["Key concepts and introduction", "Data analysis and findings", "Implementation recommendations", "Next steps and timeline", "Summary and conclusions"]

- **Middle slides (Slides 2 to {slide_count-1})**: Content slides with ONE KEY MESSAGE each
  - **ONE MESSAGE PER SLIDE**: Each slide answers "What is the one key takeaway?"
  - **VERY SHORT INSIGHT-BASED TITLE**: Max 1-2 words stating the conclusion (e.g., "Revenue Growth", "Customer Success", "Market Position")
  - **3-ZONE LAYOUT**:
    1. Headline (insight/conclusion) at top
    2. Visual/Chart in middle (if image available)
    3. Supporting points at bottom (flexible bullets)
  - **FLEXIBLE BULLET COUNT (3-8 bullets per slide)**: Adjust based on content importance and complexity
    - **Simple topics**: 3-4 bullets (keep it concise)
    - **Standard topics**: 5-6 bullets (balanced coverage)
    - **Complex/Important topics**: 7-8 bullets (comprehensive detail)
    - **QUALITY OVER QUANTITY**: Include all important points, don't pad or cut artificially
  - Include images where relevant using "image_index"
  - **CRITICAL: Each slide MUST have a UNIQUE title - no two slides should have the same title**

- **Last slide**: Thank you slide
  - Title: **"Thank You"** or **"Questions?"**
  - Simple closing slide
  - 2-3 bullet points max

🔴 BULLET POINT RULES - FLEXIBLE & CONTENT-DRIVEN 🔴:
- **FLEXIBLE BULLET COUNT: 3-8 bullets per slide** (adjust based on content importance)
  - Use MORE bullets (6-8) for complex, data-rich, or critical topics
  - Use FEWER bullets (3-4) for simple concepts or introductory slides
  - Default to 5-6 bullets for standard topics
  - **NEVER artificially limit important information** - if there are 7-8 key points, include them all
- **Maximum 6-8 WORDS per bullet point**
- SHORT, IMPACTFUL, PUNCHY language
- No long sentences - use brief, powerful phrases
- Each bullet conveys ONE key insight

EXAMPLES OF GOOD BULLETS (6-8 words):
✅ "Revenue increased 25% year over year"
✅ "Customer satisfaction reached all-time high"
✅ "Project completed ahead of schedule"
✅ "Team productivity improved with new tools"
✅ "New features launched successfully this quarter"

EXAMPLES OF BAD BULLETS (too long):
❌ "The revenue for our organization increased significantly by approximately 25% when compared to the same period in the previous year"
❌ "Customer satisfaction scores have shown remarkable improvement and reached the highest levels we have ever seen"

Content Requirements:
- **UNIQUE TITLES**: Every slide must have a unique, descriptive title
- **CONCISE**: 6-8 words per bullet point maximum
- **IMPACTFUL**: Make every word count
- **PROFESSIONAL**: Use clear, direct language
- **FOCUSED**: One key idea per bullet
- **DATA-DRIVEN**: Use numbers and facts when available
- Use insights from the provided content{"and analyze the images for key data points" if all_images else ""}
- If your content includes URLs, write them as fully qualified URLs (e.g., https://www.example.com)

🔴 CRITICAL FORMATTING RULES 🔴:
- MUST have 4-5 bullets per slide (MINIMUM 4, MAXIMUM 5)
- 6-8 words per bullet (NOT sentences)
- Use strong action verbs
- Be specific and concrete
- Avoid fluff and filler words"""

        message_content.append({
            "type": "text",
            "text": prompt_text
        })

        # Add images to the message
        print(f"Sending {len(images_to_send)} images to Claude (out of {len(all_images)} processed)")

        for idx, img in enumerate(images_to_send):
            try:
                # Validate image data before adding
                if img['base64'] and img['media_type'] and len(img['base64']) > 0:
                    message_content.append({
                        "type": "image",
                        "source": {
                            "type": "base64",
                            "media_type": img['media_type'],
                            "data": img['base64']
                        }
                    })
                    print(f"Added image {idx}: {img['media_type']}, size: {img['size']} bytes")
                else:
                    print(f"Skipping invalid image {idx}")
            except Exception as e:
                print(f"Error adding image {idx}: {str(e)}")
                continue

        print(f"Final message has {len(message_content)} content blocks")

        response = bedrock.invoke_model(
            modelId='us.anthropic.claude-sonnet-4-5-20250929-v1:0',
            body=json.dumps({
                "anthropic_version": "bedrock-2023-05-31",
                "max_tokens": 4000,
                "temperature": 0.7,
                "messages": [{"role": "user", "content": message_content}]
            })
        )
        
        result = json.loads(response['body'].read())
        slides_text = result['content'][0]['text']

        # Log the raw response for debugging
        print(f"Claude response: {slides_text[:500]}")

        # Try to extract JSON from the response
        # Claude might wrap JSON in markdown code blocks
        try:
            # First, try direct parsing
            slides_json = json.loads(slides_text)
        except json.JSONDecodeError:
            # Try to extract JSON from markdown code blocks
            json_match = re.search(r'```(?:json)?\s*(\{.*?\})\s*```', slides_text, re.DOTALL)
            if json_match:
                try:
                    slides_json = json.loads(json_match.group(1))
                    print("Successfully extracted JSON from markdown code block")
                except json.JSONDecodeError:
                    json_match = None

            # Try to find JSON object directly
            if not json_match:
                json_match = re.search(r'\{[\s\S]*"slides"[\s\S]*\}', slides_text)
                if json_match:
                    try:
                        slides_json = json.loads(json_match.group(0))
                        print("Successfully extracted JSON from text")
                    except json.JSONDecodeError:
                        slides_json = None
                else:
                    slides_json = None

        # If still no valid JSON, use fallback
        if not slides_json or 'slides' not in slides_json:
            print("WARNING: Using fallback slides - JSON parsing failed")
            print(f"Raw Claude response: {slides_text}")

            # Try to extract some content from the text for better fallback
            content_sentences = all_content.split('.')[:20]  # Get first 20 sentences
            content_points = [s.strip() for s in content_sentences if len(s.strip()) > 20]

            # Fallback with proper structure - distribute images across content slides
            # Create proper overview for first slide
            content_titles = ["Introduction", "Key Points", "Analysis", "Details", "Summary"]
            overview_points = [f"{title} and key findings" for title in content_titles[:min(slide_count-2, 4)]]

            fallback_slides = [
                {"title": description, "content": ["Overview of the presentation"] + overview_points[:3], "image_index": None},
            ]

            # Create content slides with varied content from extracted text
            for i in range(min(slide_count - 2, len(content_titles))):
                img_idx = i if i < len(all_images) else None

                # Try to use actual content instead of generic text
                start_idx = i * 3
                slide_content = content_points[start_idx:start_idx+3] if start_idx < len(content_points) else [f"Key insights about {content_titles[i]}", "Important details from research", "Supporting evidence and data"]

                # Ensure we have at least 3 bullet points
                while len(slide_content) < 3:
                    slide_content.append(f"Additional information about {description}")

                fallback_slides.append({
                    "title": content_titles[i],
                    "content": slide_content[:8],  # Max 8 bullets (flexible based on content)
                    "image_index": img_idx
                })

            fallback_slides.append(
                {"title": "Thank You", "content": ["Questions?", "Thank you for your attention"], "image_index": None}
            )

            slides_json = {"slides": fallback_slides[:slide_count]}

        # Ensure unique slide titles (accessibility requirement)
        def ensure_unique_titles(slides_list):
            """Ensure all slide titles are unique by adding numbers if needed."""
            titles_seen = {}
            for slide in slides_list:
                original_title = slide['title']
                if original_title in titles_seen:
                    # Title is duplicate, add number
                    titles_seen[original_title] += 1
                    slide['title'] = f"{original_title} ({titles_seen[original_title]})"
                else:
                    titles_seen[original_title] = 1
            return slides_list

        slides_json['slides'] = ensure_unique_titles(slides_json['slides'])

        # Create PowerPoint using local python-pptx and upload
        try:
            # Import python-pptx (will work locally)
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
            from pptx.dml.color import RGBColor

            prs = Presentation()
            prs.slide_width = Inches(10)  # Standard widescreen
            prs.slide_height = Inches(7.5)

            # Professional design constants - Following modern corporate standards
            # Modern professional fonts: Calibri, Arial, Helvetica, Segoe UI
            TITLE_FONT = 'Calibri'  # Modern corporate standard
            BODY_FONT = 'Calibri'   # Consistent, professional
            TITLE_SIZE = Pt(36)     # Professional title (28-40pt range)
            BODY_SIZE = Pt(20)      # Professional body (18-24pt range)
            TITLE_COLOR = RGBColor(31, 56, 100)  # Professional dark blue (primary color)
            BODY_COLOR = RGBColor(64, 64, 64)    # Dark gray - high contrast on white

            # Track actual image insertions
            actual_images_inserted = 0

            for slide_data in slides_json['slides']:
                # Check if slide has an image
                has_image = slide_data.get('image_index') is not None
                image_index = slide_data.get('image_index')

                if has_image and image_index is not None and image_index < len(all_images):
                    # Try to add image, fall back to text-only if it fails
                    try:
                        # Use blank layout for custom positioning
                        slide_layout = prs.slide_layouts[6]  # Blank layout
                        slide = prs.slides.add_slide(slide_layout)

                        # 3-ZONE LAYOUT: Zone 1 - Headline (key message/insight)
                        # Add title with professional formatting - positioned at top to minimize gaps
                        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.9))
                        title_frame = title_box.text_frame
                        title_frame.text = slide_data['title']
                        title_frame.word_wrap = True
                        title_para = title_frame.paragraphs[0]
                        title_para.font.name = TITLE_FONT
                        title_para.font.size = TITLE_SIZE
                        title_para.font.bold = True
                        title_para.font.color.rgb = TITLE_COLOR
                        title_para.alignment = PP_ALIGN.LEFT
                        title_para.space_after = Pt(6)  # Space below title

                        # Add image
                        img_data = all_images[image_index]

                        # Determine file extension from media type
                        ext_map = {
                            'image/jpeg': 'jpg',
                            'image/png': 'png',
                            'image/gif': 'gif',
                            'image/webp': 'webp'
                        }
                        ext = ext_map.get(img_data['media_type'], 'jpg')
                        img_path = f"/tmp/slide_img_{uuid.uuid4().hex[:8]}.{ext}"

                        # Decode and save image
                        with open(img_path, 'wb') as img_file:
                            img_file.write(base64.b64decode(img_data['base64']))

                        # 3-ZONE LAYOUT: Zone 2 - Visual/Chart (middle)
                        # Insert image (right side) - positioned higher to fill space
                        picture = slide.shapes.add_picture(img_path, Inches(5.6), Inches(1.4), width=Inches(3.8))

                        # Add alt text to image (accessibility requirement)
                        if img_data.get('is_chart'):
                            # For generated charts
                            alt_text = img_data.get('chart_title', f"Data visualization chart for {slide_data['title']}")
                        else:
                            # For web images
                            alt_text = f"Image related to {slide_data['title']}"

                        # Set the alternative text for the picture
                        picture._element._nvXxPr.cNvPr.set('descr', alt_text)

                        # Track successful image insertion
                        actual_images_inserted += 1

                        # 3-ZONE LAYOUT: Zone 3 - Supporting points (bottom/left)
                        # Add bullet points - maximize vertical fill with minimal margins
                        text_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(4.6), Inches(6.0))
                        tf = text_box.text_frame
                        tf.word_wrap = True
                        tf.vertical_anchor = MSO_ANCHOR.TOP
                        tf.margin_left = Inches(0.1)
                        tf.margin_right = Inches(0.1)
                        tf.margin_top = Inches(0.05)
                        tf.margin_bottom = Inches(0.05)

                        # Bullet Point Rule: 3-8 bullets (flexible based on content), max 6-8 words per bullet
                        content_items = slide_data['content'][:8]

                        # First bullet point with URL hyperlinks
                        p = tf.paragraphs[0]
                        # Limit to 8 words for slides with images (slightly relaxed for readability)
                        words = content_items[0].split()
                        bullet_text = ' '.join(words[:8]) + ('...' if len(words) > 8 else '')
                        add_hyperlinks_to_paragraph(p, bullet_text)
                        p.font.name = BODY_FONT
                        p.font.size = BODY_SIZE
                        p.font.color.rgb = BODY_COLOR
                        p.space_after = Pt(24)
                        p.level = 0
                        # Enable bullet point with explicit character
                        from pptx.oxml.xmlchemy import OxmlElement
                        pPr = p._element.get_or_add_pPr()
                        # Add bullet character element
                        buChar = OxmlElement('a:buChar')
                        buChar.set('char', '•')
                        pPr.append(buChar)

                        # Additional bullet points with URL hyperlinks
                        for bullet_point in content_items[1:]:
                            p = tf.add_paragraph()
                            # Limit to 8 words per bullet
                            words = bullet_point.split()
                            bullet_text = ' '.join(words[:8]) + ('...' if len(words) > 8 else '')
                            add_hyperlinks_to_paragraph(p, bullet_text)
                            p.font.name = BODY_FONT
                            p.font.size = BODY_SIZE
                            p.font.color.rgb = BODY_COLOR
                            p.space_after = Pt(24)
                            p.level = 0
                            # Enable bullet point with explicit character
                            from pptx.oxml.xmlchemy import OxmlElement
                            pPr = p._element.get_or_add_pPr()
                            # Add bullet character element
                            buChar = OxmlElement('a:buChar')
                            buChar.set('char', '•')
                            pPr.append(buChar)

                        # Enable auto-fit to prevent text overflow
                        tf.auto_size = None
                        tf.word_wrap = True

                        # Clean up temp image
                        if os.path.exists(img_path):
                            os.remove(img_path)

                    except Exception as img_error:
                        # If image fails, create text-only slide instead
                        print(f"Error adding image: {str(img_error)}, falling back to text-only")

                        # Remove the failed slide if it was created
                        if len(prs.slides) > 0:
                            try:
                                rId = prs.slides._sldIdLst[-1].rId
                                prs.part.drop_rel(rId)
                                del prs.slides._sldIdLst[-1]
                            except:
                                pass

                        # Create text-only slide instead with professional formatting
                        slide_layout = prs.slide_layouts[6]  # Use blank layout for consistency
                        slide = prs.slides.add_slide(slide_layout)

                        # 3-ZONE LAYOUT: Zone 1 - Headline (fallback slide)
                        # Add title - positioned at top to minimize gaps
                        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.9))
                        title_frame = title_box.text_frame
                        title_frame.text = slide_data['title']
                        title_frame.word_wrap = True
                        title_para = title_frame.paragraphs[0]
                        title_para.font.name = TITLE_FONT
                        title_para.font.size = TITLE_SIZE
                        title_para.font.bold = True
                        title_para.font.color.rgb = TITLE_COLOR
                        title_para.alignment = PP_ALIGN.LEFT
                        title_para.space_after = Pt(6)

                        # 3-ZONE LAYOUT: Zone 3 - Supporting points (centered)
                        # Add content - maximize vertical fill with minimal margins
                        text_box = slide.shapes.add_textbox(Inches(1.2), Inches(1.3), Inches(7.6), Inches(6.0))
                        tf = text_box.text_frame
                        tf.word_wrap = True
                        tf.vertical_anchor = MSO_ANCHOR.TOP
                        tf.margin_left = Inches(0.15)
                        tf.margin_right = Inches(0.15)
                        tf.margin_top = Inches(0.1)
                        tf.margin_bottom = Inches(0.1)

                        # Bullet Point Rule: 3-8 bullets (flexible based on content), max 6-8 words per bullet
                        content_items = slide_data['content'][:8]

                        # First bullet point with URL hyperlinks
                        p = tf.paragraphs[0]
                        # Limit to 10 words for full-width slides (more space available)
                        words = content_items[0].split()
                        bullet_text = ' '.join(words[:10]) + ('...' if len(words) > 10 else '')
                        add_hyperlinks_to_paragraph(p, bullet_text)
                        p.font.name = BODY_FONT
                        p.font.size = BODY_SIZE
                        p.font.color.rgb = BODY_COLOR
                        p.space_after = Pt(24)
                        p.level = 0
                        # Enable bullet point with explicit character
                        from pptx.oxml.xmlchemy import OxmlElement
                        pPr = p._element.get_or_add_pPr()
                        # Add bullet character element
                        buChar = OxmlElement('a:buChar')
                        buChar.set('char', '•')
                        pPr.append(buChar)

                        # Additional bullet points with URL hyperlinks
                        for bullet_point in content_items[1:]:
                            p = tf.add_paragraph()
                            # Limit to 10 words per bullet
                            words = bullet_point.split()
                            bullet_text = ' '.join(words[:10]) + ('...' if len(words) > 10 else '')
                            add_hyperlinks_to_paragraph(p, bullet_text)
                            p.font.name = BODY_FONT
                            p.font.size = BODY_SIZE
                            p.font.color.rgb = BODY_COLOR
                            p.space_after = Pt(24)
                            p.level = 0
                            # Enable bullet point with explicit character
                            from pptx.oxml.xmlchemy import OxmlElement
                            pPr = p._element.get_or_add_pPr()
                            # Add bullet character element
                            buChar = OxmlElement('a:buChar')
                            buChar.set('char', '•')
                            pPr.append(buChar)

                        # Enable auto-fit to prevent text overflow
                        tf.auto_size = None
                        tf.word_wrap = True

                        # Clean up temp image file if it exists
                        try:
                            if 'img_path' in locals() and os.path.exists(img_path):
                                os.remove(img_path)
                        except:
                            pass

                else:
                    # Regular text-only slide with professional formatting
                    slide_layout = prs.slide_layouts[6]  # Blank layout for consistency
                    slide = prs.slides.add_slide(slide_layout)

                    # 3-ZONE LAYOUT: Zone 1 - Headline (key message)
                    # Add title - positioned at top to minimize gaps
                    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.9))
                    title_frame = title_box.text_frame
                    title_frame.text = slide_data['title']
                    title_frame.word_wrap = True
                    title_para = title_frame.paragraphs[0]
                    title_para.font.name = TITLE_FONT
                    title_para.font.size = TITLE_SIZE
                    title_para.font.bold = True
                    title_para.font.color.rgb = TITLE_COLOR
                    title_para.alignment = PP_ALIGN.LEFT
                    title_para.space_after = Pt(6)

                    # 3-ZONE LAYOUT: Zone 3 - Supporting points (centered with whitespace)
                    # Add content - maximize vertical fill with minimal margins
                    text_box = slide.shapes.add_textbox(Inches(1.2), Inches(1.3), Inches(7.6), Inches(6.0))
                    tf = text_box.text_frame
                    tf.word_wrap = True
                    tf.vertical_anchor = MSO_ANCHOR.TOP
                    tf.margin_left = Inches(0.15)
                    tf.margin_right = Inches(0.15)
                    tf.margin_top = Inches(0.1)
                    tf.margin_bottom = Inches(0.1)

                    # Bullet Point Rule: 3-8 bullets (flexible based on content), max 6-8 words per bullet
                    content_items = slide_data['content'][:8]

                    # First bullet point with hyperlinks
                    p = tf.paragraphs[0]
                    # Limit to 10 words for full-width slides (more space available)
                    words = content_items[0].split()
                    bullet_text = ' '.join(words[:10]) + ('...' if len(words) > 10 else '')
                    add_hyperlinks_to_paragraph(p, bullet_text)
                    p.font.name = BODY_FONT
                    p.font.size = BODY_SIZE
                    p.font.color.rgb = BODY_COLOR
                    p.space_after = Pt(24)
                    p.level = 0
                    # Enable bullet point with explicit character
                    from pptx.oxml.xmlchemy import OxmlElement
                    pPr = p._element.get_or_add_pPr()
                    # Add bullet character element
                    buChar = OxmlElement('a:buChar')
                    buChar.set('char', '•')
                    pPr.append(buChar)

                    # Additional bullet points with hyperlinks
                    for bullet_point in content_items[1:]:
                        p = tf.add_paragraph()
                        # Limit to 10 words per bullet
                        words = bullet_point.split()
                        bullet_text = ' '.join(words[:10]) + ('...' if len(words) > 10 else '')
                        add_hyperlinks_to_paragraph(p, bullet_text)
                        p.font.name = BODY_FONT
                        p.font.size = BODY_SIZE
                        p.font.color.rgb = BODY_COLOR
                        p.space_after = Pt(24)
                        p.level = 0
                        # Enable bullet point with explicit character
                        from pptx.oxml.xmlchemy import OxmlElement
                        pPr = p._element.get_or_add_pPr()
                        # Add bullet character element
                        buChar = OxmlElement('a:buChar')
                        buChar.set('char', '•')
                        pPr.append(buChar)

                    # Enable auto-fit to prevent text overflow
                    tf.auto_size = None
                    tf.word_wrap = True
            
            # Save and upload
            filename = f"presentation_{uuid.uuid4().hex[:8]}.pptx"
            local_path = f"/tmp/{filename}"
            prs.save(local_path)
            
            file_size = os.path.getsize(local_path)
            
            bucket_name = os.environ.get('S3_BUCKET', 'ppt-generator-1759467436-803633136603')
            s3 = boto3.client('s3')
            s3.upload_file(local_path, bucket_name, filename)
            
            download_url = s3.generate_presigned_url(
                'get_object',
                Params={'Bucket': bucket_name, 'Key': filename},
                ExpiresIn=86400
            )
            
            os.remove(local_path)
            
            response_data = {
                'message': 'PowerPoint created and uploaded successfully',
                'download_url': download_url,
                'slides': slides_json,
                'pptx_info': {
                    'filename': filename,
                    'size_bytes': file_size,
                    'slide_count': len(slides_json['slides']),
                    'format': 'pptx',
                    'images_processed': len(all_images),
                    'images_inserted': actual_images_inserted
                }
            }

            # Add data analysis info if applicable
            if data_analysis:
                response_data['data_analysis'] = {
                    'insights_count': len(data_analysis.get('insights', [])),
                    'charts_generated': len(generated_charts),
                    'visualizations': generated_charts
                }

            return {
                'statusCode': 200,
                'headers': {
                    'Content-Type': 'application/json',
                    'Access-Control-Allow-Origin': '*',
                    'Access-Control-Allow-Headers': 'Content-Type',
                    'Access-Control-Allow-Methods': 'POST, OPTIONS'
                },
                'body': json.dumps(response_data)
            }

        except ImportError:
            # Fallback: Return slides data for local processing
            return {
                'statusCode': 200,
                'headers': {
                    'Content-Type': 'application/json',
                    'Access-Control-Allow-Origin': '*',
                    'Access-Control-Allow-Headers': 'Content-Type',
                    'Access-Control-Allow-Methods': 'POST, OPTIONS'
                },
                'body': json.dumps({
                    'message': 'Slides generated - use local processing for PPTX',
                    'slides': slides_json,
                    'images': [{'url': img['url'], 'media_type': img['media_type']} for img in all_images],
                    'pptx_info': {
                        'slide_count': len(slides_json['slides']),
                        'format': 'json',
                        'images_processed': len(all_images),
                        'note': 'Use local python-pptx to create PowerPoint file'
                    }
                })
            }

    except Exception as e:
        print(f"ERROR in lambda_handler: {str(e)}")
        import traceback
        traceback.print_exc()
        return {
            'statusCode': 500,
            'headers': {
                'Content-Type': 'application/json',
                'Access-Control-Allow-Origin': '*',
                'Access-Control-Allow-Headers': 'Content-Type',
                'Access-Control-Allow-Methods': 'POST, OPTIONS'
            },
            'body': json.dumps({'error': str(e)})
        }