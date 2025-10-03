from pptx import Presentation
import streamlit as st
import tempfile
import os

def create_pptx_from_slides(slides_data):
    """Create PPTX file from slides data"""
    prs = Presentation()
    
    for i, slide_data in enumerate(slides_data['slides']):
        # Use title slide layout for first slide, content layout for others
        if i == 0:
            slide_layout = prs.slide_layouts[0]  # Title slide layout
            slide = prs.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            title_shape.text = slide_data['title']
            if slide.shapes.placeholders[1]:  # Subtitle
                subtitle_shape = slide.shapes.placeholders[1]
                subtitle_shape.text = "\n".join(slide_data['content'])
        else:
            slide_layout = prs.slide_layouts[1]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            title_shape = slide.shapes.title
            title_shape.text = slide_data['title']
            
            # Add bullet points
            content = slide.placeholders[1]
            tf = content.text_frame
            tf.text = slide_data['content'][0]
            
            for bullet_point in slide_data['content'][1:]:
                p = tf.add_paragraph()
                p.text = bullet_point
                p.level = 0
    
    # Save to temporary file
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
    prs.save(temp_file.name)
    return temp_file.name